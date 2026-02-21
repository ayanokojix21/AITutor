"""Generate Eduverse hackathon pitch deck (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy
ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # Purple
ACCENT2   = RGBColor(0x00, 0xD2, 0xFF)   # Cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xB8, 0xCC)   # Muted text
GREEN     = RGBColor(0x00, 0xE6, 0x76)   # Success green
ORANGE    = RGBColor(0xFF, 0x9F, 0x43)   # Warm accent

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text,
             size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_slide(slide, left, top, width, height, items, size=16, color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return tf


def accent_bar(slide, left, top, width=Inches(0.8), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card(slide, left, top, width, height, title, body, icon="", title_color=ACCENT2):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x3B)
    shape.line.fill.background()
    # Title
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
             f"{icon}  {title}", size=16, color=title_color, bold=True)
    # Body
    add_text(slide, left + Inches(0.3), top + Inches(0.65), width - Inches(0.6), height - Inches(0.85),
             body, size=13, color=LIGHT)


def section_number(slide, num, label):
    add_text(slide, Inches(0.8), Inches(0.4), Inches(1), Inches(0.6),
             f"{num:02d}", size=36, color=ACCENT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(0.6))
    add_text(slide, Inches(0.8), Inches(1.15), Inches(4), Inches(0.5),
             label.upper(), size=12, color=LIGHT, bold=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative accent
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(5), Inches(5))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.fill.fore_color.brightness = -0.7
shape.line.fill.background()

add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
         "EDUVERSE", size=56, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.5), Inches(9), Inches(0.8),
         "The Next-Gen Multimodal AI Tutor", size=28, color=ACCENT2)
accent_bar(slide, Inches(1), Inches(3.4), Inches(1.5), ACCENT)
add_text(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.6),
         "Transforming passive video content into active, intelligent knowledge bases.",
         size=16, color=LIGHT)
add_text(slide, Inches(1), Inches(6.0), Inches(8), Inches(0.4),
         "Tech Stack: LangGraph  |  Llama-3 (Groq)  |  Supabase pgvector  |  FastAPI",
         size=12, color=RGBColor(0x80, 0x88, 0x99))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 1, "Problem Statement")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         '"Educational content is rich — but opaque."', size=28, color=WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(2.8), Inches(3.8), Inches(3.5),
         "Information Overload", "Students spend hours scrubbing through video timelines to find a single concept. 72% report video-search as a major pain point.", "📚")
add_card(slide, Inches(4.9), Inches(2.8), Inches(3.8), Inches(3.5),
         "Modality Blindness", "Standard search only indexes text. It misses spoken explanations in video, diagrams in slides, and whiteboard writing.", "🔇", ORANGE)
add_card(slide, Inches(9.0), Inches(2.8), Inches(3.8), Inches(3.5),
         "Unscalable Support", "1:1 tutoring is expensive. AI chatbots hallucinate because they lack course context.", "💰", GREEN)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — MOTIVATION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 2, "Motivation")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Why Now? The Trust Gap", size=28, color=WHITE, bold=True)

# Problem vs Solution columns
add_text(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.4), "THE PROBLEM", size=14, color=ACCENT, bold=True)
add_bullet_slide(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2),
    ["• Generic LLMs hallucinate facts", "• Students need answers they can TRUST", "• 'Trust but verify' is impossible without sources"], size=15)

add_text(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(0.4), "THE SOLUTION: CITATIONS", size=14, color=GREEN, bold=True)
add_bullet_slide(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(2),
    ["• Every answer must cite its source (Timestamp/Page)", "• Evidence-based RAG architecture", "• Zero Hallucination goal"], size=15)

add_card(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2),
         "💡 Key Insight", "The value isn't just the answer — it's the EVIDENCE. Eduverse provides Citation-Specific Answers.", "", ACCENT2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — APPLICATION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 3, "Application")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Real-World Use Cases", size=28, color=WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(2.7), Inches(3.8), Inches(2.2), "🎓 Universities", "Automated TAs. 'Where did the professor discuss Bias Variance?' → In Lecture 5 at 42:10.", "📍", ACCENT2)
add_card(slide, Inches(4.9), Inches(2.7), Inches(3.8), Inches(2.2), "⚖️ Legal Tech", "Citation-specific search. 'Find the exact moment the witness contradicted the deposition.'", "⚖️", ACCENT2)
add_card(slide, Inches(9.0), Inches(2.7), Inches(3.8), Inches(2.2), "⚕️ Medical", "Procedure analysis. 'Show step-by-step suture technique from training video.'", "⚕️", ACCENT2)

add_text(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.4), "CORE CAPABILITIES", size=14, color=ACCENT, bold=True)
add_bullet_slide(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1.5),
    ["✦ Citation-Specific Retrieval — Exact timestamps for videos", "✦ Semantic Normalizer — Aligns spoken word with formal text", "✦ Temporal Understanding — AI knows the sequence of concepts"], size=14)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — PROPOSED METHOD (Architecture)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 4, "Proposed Method")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Multimodal RAG + Agentic Workflows", size=28, color=WHITE, bold=True)

# Architecture flow boxes
stages = [
    ("📥 INGEST", "Drive Sync\nClassroom API", ACCENT, Inches(0.5)),
    ("🔄 PROCESS", "Whisper (Audio)\nCV (Frames)", ACCENT2, Inches(2.9)),
    ("⚖️ NORMALIZE", "Semantic\nRestructuring", ORANGE, Inches(5.3)),
    ("🧠 EMBED", "pgvector\nHybrid Search", GREEN, Inches(7.7)),
    ("💬 CITATION", "Retrieval with\nTimestamps", RGBColor(0xFF, 0x63, 0x84), Inches(10.1)),
]
for title, body, color, left in stages:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.7), Inches(2.2), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x3B)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    add_text(slide, left+Inches(0.1), Inches(2.8), Inches(2.0), Inches(0.4), title, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left+Inches(0.1), Inches(3.2), Inches(2.0), Inches(1.2), body, size=11, color=LIGHT, align=PP_ALIGN.CENTER)

# Arrows
for left in [Inches(2.7), Inches(5.1), Inches(7.5), Inches(9.9)]:
    add_text(slide, left, Inches(3.3), Inches(0.3), Inches(0.5), "→", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Tech cards logic kept conceptually similar but shortened for fit


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — SEMANTIC NORMALIZER
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 5, "Semantic Normalizer")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Bridging the Modality Gap", size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.4), "THE CHALLENGE", size=14, color=RGBColor(0xFF, 0x6B, 0x6B), bold=True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(1.5),
         "Transcribed speech is messy. Textbooks are formal. Latent Space Mismatch: 'Gradient Descent' verbally looks different from a textbook definition.", size=14, color=LIGHT)

add_text(slide, Inches(7.0), Inches(2.6), Inches(5), Inches(0.4), "THE SOLUTION", size=14, color=GREEN, bold=True)
add_text(slide, Inches(7.0), Inches(3.0), Inches(5), Inches(1.5),
         "A pre-processing layer that 'normalizes' all content to a unified semantic style BEFORE embedding.", size=14, color=LIGHT)

# Diagram
y = 4.8
add_card(slide, Inches(0.8), Inches(y), Inches(3.5), Inches(1.8), "Raw Transcript", "'Um, so basically if you go down the hill...'", "🗣️", ACCENT2)
add_text(slide, Inches(4.5), Inches(y+0.5), Inches(0.5), Inches(0.5), "→", size=24, color=WHITE, bold=True)
add_card(slide, Inches(5.2), Inches(y-0.2), Inches(3.0), Inches(2.2), "Semantic Normalizer", "LLM Rewriter + Context Injection", "⚙️", ORANGE)
add_text(slide, Inches(8.4), Inches(y+0.5), Inches(0.5), Inches(0.5), "→", size=24, color=WHITE, bold=True)
add_card(slide, Inches(9.1), Inches(y), Inches(3.5), Inches(1.8), "Normalized Chunk", "'Gradient descent optimizes parameters by moving steepest descent...'", "📄", GREEN)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — DATASETS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 6, "Datasets & Privacy")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Bring Your Own Data (BYOD)", size=28, color=WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(1.5), "📂 Supported Formats", "PDF, MP4, MP3, PPTX, Google Slides. System matches diverse input types.", "", ACCENT2)
add_card(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(1.5), "🔗 Ingestion", "Google Classroom Sync, Drive Integration, Direct Upload.", "", ACCENT)
add_card(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.2), "🔒 Privacy", "• No training on user data\n• Private vector namespaces\n• Encrypted credentials", "", GREEN)
add_card(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.2), "📊 Benchmarks", "• MIT OpenCourseWare\n• Khan Academy Transcripts\n• Custom Q&A Test Suite", "", ORANGE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — EXPERIMENTS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 7, "Experiments & Validation")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "Measuring Success", size=28, color=WHITE, bold=True)

metrics = [
    ("Retrieval Recall@5", "> 90%", "Finding the correct video segment.", "🎯", ACCENT2),
    ("Response Latency", "< 500ms", "Real-time feel via Groq LPUs.", "⚡", GREEN),
    ("Hallucination Rate", "< 5%", "Verified against ground truth.", "🛡️", ORANGE),
    ("User Satisfaction", "8.5/10", "System Usability Scale (SUS).", "⭐", ACCENT),
]
for i, (t, m, d, ic, c) in enumerate(metrics):
    left = Inches(0.8 + i*3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.6), Inches(2.8), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x3B)
    shape.line.color.rgb = c
    shape.line.width = Pt(2)
    add_text(slide, left+Inches(0.2), Inches(2.75), Inches(2.4), Inches(0.4), f"{ic}  {t}", size=14, color=c, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left+Inches(0.2), Inches(3.3), Inches(2.4), Inches(0.6), m, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left+Inches(0.2), Inches(4.1), Inches(2.4), Inches(1.8), d, size=12, color=LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — NOVELTY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_number(slide, 8, "Novelty & Scope to Scale")

add_text(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.7),
         "What Makes Eduverse Unique", size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.4), "🔬 NOVELTY", size=16, color=ACCENT2, bold=True)
novelty = [
    "True Multimodality: Linking Audio + Visual + Text",
    "Citation-Specific Accuracy: Zero Hallucination",
    "Temporal Awareness: Knows 'when' concepts happen",
    "Semantic Normalizer: Better vector alignment"
]
for i, item in enumerate(novelty):
    add_card(slide, Inches(0.8), Inches(3.1 + i*0.95), Inches(5.5), Inches(0.8), f"{'▸'}", item, "", ACCENT2)

add_text(slide, Inches(7.0), Inches(2.6), Inches(5.5), Inches(0.4), "📈 SCALABILITY", size=16, color=GREEN, bold=True)
scale = [
    ("Education_Legal", "Case law video analysis"),
    ("Education_Medical", "Surgical procedure training"),
    ("Enterprise_Ready", "Supabase serverless scale"),
    ("Global_Reach", "Whisper 90+ languages")
]
for i, (t, d) in enumerate(scale):
    add_card(slide, Inches(7.0), Inches(3.1 + i*0.95), Inches(5.5), Inches(0.8), t.replace("_", " → "), d, "", GREEN)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(3), Inches(6), Inches(6))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.fill.fore_color.brightness = -0.7
shape.line.fill.background()

add_text(slide, Inches(0), Inches(2.0), Inches(13.3), Inches(1), "EDUVERSE", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.3), Inches(13.3), Inches(0.8), "Transforming how the world learns — one lecture at a time.", size=24, color=ACCENT2, align=PP_ALIGN.CENTER)
accent_bar(slide, Inches(5.8), Inches(4.3), Inches(1.8), ACCENT)
add_text(slide, Inches(0), Inches(5.0), Inches(13.3), Inches(1), "Not just a chatbot. A cognitive engine.", size=20, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(6.5), Inches(13.3), Inches(0.5), "Thank You", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# SAVE
output = r"c:\Users\HP\et-genai\AITutor\Eduverse_GenAI_Pitch.pptx"
prs.save(output)
print(f"✅ Saved 10-slide deck to: {output}")
